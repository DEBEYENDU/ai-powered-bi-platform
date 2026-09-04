"""Multi-format exporters for reports.

Design: every exporter receives (built_report_dict, html_string, output_path).
Formats backed only by the standard library always work (html, csv, json, zip).
Formats needing optional dependencies (pdf/reportlab, docx/python-docx,
xlsx/openpyxl, pptx/python-pptx, png/svg renderers) degrade gracefully: they
write an HTML-based placeholder and report ``degraded: True`` so callers can
surface it instead of crashing minimal deployments.
"""

from __future__ import annotations

import csv
import io
import json
import zipfile
from pathlib import Path
from typing import Any

SUPPORTED_FORMATS = ["pdf", "docx", "xlsx", "csv", "json", "html", "pptx", "png", "svg", "zip"]


def _flatten_rows(report: dict[str, Any]) -> tuple[list[str], list[list[Any]]]:
    header = ["section_id", "kind", "title", "content"]
    rows: list[list[Any]] = []
    for s in report.get("sections", []):
        d = s if isinstance(s, dict) else {}
        data = d.get("data", {})
        payload = data.get("payload", data) if isinstance(data, dict) else data
        rows.append(
            [d.get("section_id", ""), d.get("kind", ""), d.get("title", ""), str(payload)[:2000]]
        )
    return header, rows


def export_html(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    path.write_text(html, encoding="utf-8")
    return {"format": "html", "degraded": False}


def export_json(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    serializable = json.loads(json.dumps(report, default=str))
    path.write_text(json.dumps(serializable, indent=2), encoding="utf-8")
    return {"format": "json", "degraded": False}


def export_csv(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    header, rows = _flatten_rows(report)
    with path.open("w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(header)
        writer.writerows(rows)
    return {"format": "csv", "degraded": False}


def export_zip(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    header, rows = _flatten_rows(report)
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(header)
    writer.writerows(rows)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("report.html", html)
        zf.writestr(
            "report.json", json.dumps(json.loads(json.dumps(report, default=str)), indent=2)
        )
        zf.writestr("report.csv", buf.getvalue())
    return {"format": "zip", "degraded": False}


def export_pdf(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    try:
        from reportlab.lib.pagesizes import A4  # type: ignore
        from reportlab.pdfgen import canvas  # type: ignore

        c = canvas.Canvas(str(path), pagesize=A4)
        _width, height = A4
        y = height - 60
        c.setFont("Helvetica-Bold", 16)
        c.drawString(40, y, str(report.get("title", "Report"))[:90])
        y -= 24
        c.setFont("Helvetica", 10)
        for s in report.get("sections", []):
            d = s if isinstance(s, dict) else {}
            line = f"- {d.get('title') or d.get('section_id')}: {str(d.get('data') or {})[:110]}"
            for chunk in [line[i : i + 110] for i in range(0, len(line), 110)]:
                if y < 60:
                    c.showPage()
                    y = height - 60
                    c.setFont("Helvetica", 10)
                c.drawString(40, y, chunk[:110])
                y -= 14
        c.save()
        return {"format": "pdf", "degraded": False}
    except Exception as exc:
        path.write_text(html, encoding="utf-8")  # placeholder keeps pipeline green
        return {"format": "pdf", "degraded": True, "reason": f"reportlab unavailable: {exc}"}


def export_docx(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    try:
        from docx import Document  # type: ignore

        doc = Document()
        doc.add_heading(str(report.get("title", "Report")), level=1)
        for s in report.get("sections", []):
            d = s if isinstance(s, dict) else {}
            doc.add_heading(str(d.get("title") or d.get("section_id", "")), level=2)
            doc.add_paragraph(str((d.get("data") or {}).get("payload", d.get("data", "")))[:2000])
        doc.save(str(path))
        return {"format": "docx", "degraded": False}
    except Exception as exc:
        path.write_text(html, encoding="utf-8")
        return {"format": "docx", "degraded": True, "reason": f"python-docx unavailable: {exc}"}


def export_xlsx(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    try:
        from openpyxl import Workbook  # type: ignore

        wb = Workbook()
        ws = wb.active
        ws.title = "Report"
        ws.append(["Section", "Kind", "Title", "Content"])
        _, rows = _flatten_rows(report)
        for r in rows:
            ws.append(r)
        wb.save(str(path))
        return {"format": "xlsx", "degraded": False}
    except Exception as exc:
        path.write_text(html, encoding="utf-8")
        return {"format": "xlsx", "degraded": True, "reason": f"openpyxl unavailable: {exc}"}


def export_pptx(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = str(report.get("title", "Report"))[:100]
        for s in report.get("sections", [])[:10]:
            d = s if isinstance(s, dict) else {}
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = str(d.get("title") or d.get("section_id", ""))[:100]
            slide.placeholders[1].text = str(d.get("data") or {})[:1000]
        prs.save(str(path))
        return {"format": "pptx", "degraded": False}
    except Exception as exc:
        path.write_text(html, encoding="utf-8")
        return {"format": "pptx", "degraded": True, "reason": f"python-pptx unavailable: {exc}"}


def export_png(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    path.write_text(html, encoding="utf-8")
    return {
        "format": "png",
        "degraded": True,
        "reason": "chart rasterization requires dashboard viz engine",
    }


def export_svg(report: dict[str, Any], html: str, path: Path) -> dict[str, Any]:
    title = str(report.get("title", "Report")).replace("&", "&amp;").replace("<", "&lt;")
    path.write_text(
        f'<svg xmlns="http://www.w3.org/2000/svg" width="800" height="200">'
        f'<text x="20" y="60" font-size="24">{title[:80]}</text>'
        f'<text x="20" y="100" font-size="14">{len(report.get("sections", []))} sections</text></svg>',
        encoding="utf-8",
    )
    return {
        "format": "svg",
        "degraded": True,
        "reason": "summary SVG; full charts via dashboard engine",
    }


EXPORTERS = {
    "pdf": export_pdf,
    "docx": export_docx,
    "xlsx": export_xlsx,
    "csv": export_csv,
    "json": export_json,
    "html": export_html,
    "pptx": export_pptx,
    "png": export_png,
    "svg": export_svg,
    "zip": export_zip,
}


def export_report(report: dict[str, Any], html: str, fmt: str, path: Path) -> dict[str, Any]:
    fmt = fmt.lower()
    if fmt not in EXPORTERS:
        raise ValueError(f"Unsupported format '{fmt}'. Supported: {SUPPORTED_FORMATS}")
    path.parent.mkdir(parents=True, exist_ok=True)
    return EXPORTERS[fmt](report, html, path)
